from pptx import Presentation
from pptx.util import Inches, Pt, Emu

pptx_path = r'c:\Users\gaoxu\.trae-cn\attachments\6a7dee59db35f9bcbe161e89\b2f215b7-2b04-4adc-9ee0-86ae4abf6090_8b792b31-2b55-4c92-a079-e43ae7d69933_智启题库-大学生创新创业大赛答辩-优化版-修正.pptx'
prs = Presentation(pptx_path)

print(f'Slide width: {prs.slide_width}, height: {prs.slide_height}')
print(f'Slide width inches: {prs.slide_width/914400:.2f}, height inches: {prs.slide_height/914400:.2f}')
print(f'Number of slides: {len(prs.slides)}')
print()

# Check slide layouts
print('=== Slide Layouts ===')
for i, layout in enumerate(prs.slide_layouts):
    print(f'Layout {i}: {layout.name}')
print()

# Inspect ALL slides for shapes and watermarks
for slide_idx, slide in enumerate(prs.slides):
    print(f'=== Slide {slide_idx+1} ===')
    for shape_idx, shape in enumerate(slide.shapes):
        text = ''
        if shape.has_text_frame:
            text = shape.text_frame.text[:100]
        shape_info = f'  Shape {shape_idx}: type={shape.shape_type}, name="{shape.name}", left={shape.left}, top={shape.top}, w={shape.width}, h={shape.height}'
        if text:
            shape_info += f', text="{text}"'
        if 'AI' in (shape.name or '') or '生成' in text:
            shape_info += ' [WATERMARK?]'
        print(shape_info)
    print()
