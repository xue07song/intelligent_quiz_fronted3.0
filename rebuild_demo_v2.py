"""
重建功能展示幻灯片 v2：
1. 裁剪截图为 4:3 横向格式（保留顶部重要内容）
2. 每张幻灯片放置 3 张截图
3. 压缩文字内容区域
"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image as PILImage

# ==================== 常量 ====================
PPTX_PATH = r'C:\Users\gaoxu\Desktop\智启题库-大学生创新创业大赛答辩-商业计划书.pptx'
SCREENSHOT_DIR = r'c:\Users\gaoxu\AppData\Local\Temp\trae\screenshots'
CROP_DIR = r'd:\intelligent_quiz_fronted3.0\cropped_shots'

NAVY = RGBColor(0x1A, 0x2B, 0x5E)
TEAL = RGBColor(0x00, 0xB8, 0xA3)
BLUE = RGBColor(0x2D, 0x6C, 0xDF)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT = RGBColor(0x9C, 0xA3, 0xAF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT_NAME = "Microsoft YaHei"

# ==================== 辅助函数 ====================
def set_font(run, size_pt, color, bold=False, name=FONT_NAME):
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = name

def add_textbox(slide, left, top, width, height, text, size_pt, color, bold=False, align=None):
    box = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_font(run, size_pt, color, bold)
    if align:
        p.alignment = align
    return box

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_compact_card(slide, left, top, width, height, accent_color, title, items, title_size=12):
    add_rect(slide, left, top, width, height, WHITE)
    add_rect(slide, left, top, 38100, height, accent_color)
    add_textbox(slide, left + 114300, top + 76200, width - 190500, 201168, title, title_size, NAVY, True)
    y = top + 304800
    for item in items:
        add_textbox(slide, left + 114300, y, width - 190500, 182880, f"•  {item}", 8, GRAY)
        y += 219456
    return y

def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)

def clear_slide_content(slide, keep_count=4):
    shapes_to_remove = []
    for i, shape in enumerate(slide.shapes):
        if i >= keep_count:
            shapes_to_remove.append(shape)
    for s in shapes_to_remove:
        remove_shape(s)

def crop_screenshot_to_landscape(src_path, dst_path, target_ratio=4/3):
    """裁剪截图顶部，保持 target_ratio 宽高比"""
    with PILImage.open(src_path) as img:
        w, h = img.size
        target_h = int(w / target_ratio)
        if target_h > h:
            target_h = h
        cropped = img.crop((0, 0, w, target_h))
        cropped.save(dst_path)
        print(f"  裁剪: {os.path.basename(src_path)} {w}x{h} -> {w}x{target_h}")

def get_image_aspect_ratio(path):
    with PILImage.open(path) as img:
        w, h = img.size
        return w / h

def fit_image_in_area(path, max_w, max_h):
    ratio = get_image_aspect_ratio(path)
    w = max_w
    h = int(w / ratio)
    if h > max_h:
        h = max_h
        w = int(h * ratio)
    return w, h

def add_screenshot_row(slide, img_specs, top, max_h, total_width, left_margin):
    """添加一行截图，自动等宽排列"""
    n = len(img_specs)
    gap = 182880  # 0.2"
    each_w = (total_width - (n - 1) * gap) // n
    x = left_margin
    for img_path, caption in img_specs:
        w, h = fit_image_in_area(img_path, each_w, max_h)
        cx = x + (each_w - w) // 2
        cy = top + (max_h - h) // 2
        slide.shapes.add_picture(img_path, Emu(cx), Emu(cy), Emu(w), Emu(h))
        add_textbox(slide, x, top + max_h + 10000, each_w, 182880, caption, 7, GRAY_LIGHT, align=PP_ALIGN.CENTER)
        x += each_w + gap

def find_screenshot(name):
    p = os.path.join(SCREENSHOT_DIR, name)
    return p if os.path.exists(p) else None

def ensure_cropped(filename, crop_ratio=4/3):
    """获取裁剪后的截图路径，如不存在则创建"""
    src = find_screenshot(filename)
    if not src:
        return None
    os.makedirs(CROP_DIR, exist_ok=True)
    dst = os.path.join(CROP_DIR, filename)
    if not os.path.exists(dst):
        crop_screenshot_to_landscape(src, dst, crop_ratio)
    return dst

# ==================== 重建函数 ====================
def rebuild_student_slide(slide):
    clear_slide_content(slide, keep_count=4)

    add_textbox(slide, 594360, 1005840, 8092440, 274320,
                "6种题型在线答题 · 自适应难度练习 · 错题归集分析 · AI智能助手", 11, GRAY)

    # 4张紧凑卡片 (2x2)
    cw = 3977640
    ch = 760000
    cx1 = 457200
    cx2 = 4659840
    cy1 = 1280160
    cy2 = 1280160 + ch + 120000  # 2160160

    add_compact_card(slide, cx1, cy1, cw, ch, TEAL, "在线答题", [
        "6种题型 · 答题卡导航 · 实时计时 · 断点续答",
        "主观题标注「语义评阅·教师可复核」",
    ])
    add_compact_card(slide, cx2, cy1, cw, ch, BLUE, "自适应练习", [
        "从1级开始，逐题动态调整难度",
        "库存预检 + 方案推荐 · AI即时答疑",
    ])
    add_compact_card(slide, cx1, cy2, cw, ch, AMBER, "错题本与分析", [
        "错题自动归集 · 一键重练",
        "章节/题型/难度三维统计 · SVG趋势图",
    ])
    add_compact_card(slide, cx2, cy2, cw, ch, NAVY, "AI智能助手", [
        "浮动球设计 · 可拖拽对话面板",
        "自然语言组卷 · 错题浓缩 · 同类题推荐",
    ])

    # 3张截图
    shots = [
        (ensure_cropped("student_home.png"), "试卷列表"),
        (ensure_cropped("student_exam.png"), "在线答题"),
        (ensure_cropped("student_practice.png"), "自适应练习"),
    ]
    valid_shots = [(p, c) for p, c in shots if p]
    add_screenshot_row(slide, valid_shots, top=3050000, max_h=1600000,
                       total_width=8092440, left_margin=594360)
    add_textbox(slide, 594360, 4750000, 8092440, 182880,
                "来源：智启题库系统实拍（localhost:5173）", 8, GRAY_LIGHT, align=PP_ALIGN.CENTER)
    print(f"  学生端: {len(valid_shots)} 张截图")

def rebuild_teacher_slide(slide):
    clear_slide_content(slide, keep_count=4)

    add_textbox(slide, 594360, 1005840, 8092440, 274320,
                "多约束智能组卷 · AI自动出题 · 试卷分析与学情总览", 11, GRAY)

    cw = 3977640
    ch = 760000
    cx1 = 457200
    cx2 = 4659840
    cy1 = 1280160
    cy2 = 2160160

    add_compact_card(slide, cx1, cy1, cw, ch, TEAL, "多约束智能组卷", [
        "四步配置 · 9套预设 + 手动自定义 + AI辅助",
        "四维校验 · 实时库存检查 · 替代方案推荐",
    ])
    add_compact_card(slide, cx2, cy1, cw, ch, BLUE, "AI自动出题", [
        "按章节/知识点/题型/难度生成草稿",
        "人工审核 → 入库 · Excel批量导入",
    ])
    add_compact_card(slide, cx1, cy2, cw, ch, AMBER, "试卷分析与学情", [
        "正确率分析 · 成绩分布 · 班级对比",
        "主观题复核 · 学生个性化分析",
    ])
    add_compact_card(slide, cx2, cy2, cw, ch, NAVY, "自适应学情总览", [
        "练习人数 · 有效次数 · 整体正确率",
        "关注学生标记 · 建议跳转组卷",
    ])

    shots = [
        (ensure_cropped("teacher_exam_gen.png"), "智能组卷"),
        (ensure_cropped("teacher_analytics.png"), "试卷分析"),
        (ensure_cropped("teacher_questions.png"), "题库管理"),
    ]
    valid_shots = [(p, c) for p, c in shots if p]
    add_screenshot_row(slide, valid_shots, top=3050000, max_h=1600000,
                       total_width=8092440, left_margin=594360)
    add_textbox(slide, 594360, 4750000, 8092440, 182880,
                "来源：智启题库系统实拍（localhost:5173）", 8, GRAY_LIGHT, align=PP_ALIGN.CENTER)
    print(f"  教师端: {len(valid_shots)} 张截图")

def rebuild_admin_slide(slide):
    clear_slide_content(slide, keep_count=4)

    add_textbox(slide, 594360, 1005840, 8092440, 274320,
                "三角色RBAC权限管理 · 注册审批闭环 · 全局数据治理", 11, GRAY)

    # 3张内容卡片
    cw = 2545600
    ch = 1280000
    cx1 = 457200
    cx2 = 3231400
    cx3 = 6005600
    cy = 1280160

    def add_admin_card(left, accent, title, items):
        add_rect(slide, left, cy, cw, ch, WHITE)
        add_rect(slide, left, cy, 38100, ch, accent)
        add_textbox(slide, left + 114300, cy + 76200, cw - 190500, 201168, title, 12, NAVY, True)
        y = cy + 304800
        for item in items:
            add_textbox(slide, left + 114300, y, cw - 190500, 182880, f"•  {item}", 8, GRAY)
            y += 219456

    add_admin_card(cx1, TEAL, "用户与权限管理", [
        "三角色RBAC：admin/teacher/student",
        "教师科目级数据隔离",
        "用户启停 · 重置密码 · 批量操作",
        "JWT 7天有效 · bcrypt加盐哈希",
    ])
    add_admin_card(cx2, BLUE, "注册审批闭环", [
        "非自助注册 · 审批制",
        "教师申请含学院/专业/科目/工号",
        "管理员审核 → 通过/拒绝",
        "通过后自动创建账号并关联科目",
    ])
    add_admin_card(cx3, AMBER, "全局数据治理", [
        "全局题库与学习数据统计",
        "增量字段兼容 · 减少迁移风险",
        "用户反馈管理：提交/回复/跟踪",
        "系统运行监控与日志",
    ])

    # 数据统计指标
    stats = [("70+", "API接口"), ("15+", "数据表"), ("30+", "前端组件"), ("6", "题型"), ("2", "大模型")]
    stat_y = 2700000
    stat_w = 1500000
    stat_total_w = 5 * stat_w + 4 * 60000
    stat_x = (SLIDE_W - stat_total_w) // 2
    for num, label in stats:
        add_textbox(slide, stat_x, stat_y, stat_w, 274320, num, 16, TEAL, True, align=PP_ALIGN.CENTER)
        add_textbox(slide, stat_x, stat_y + 274320, stat_w, 182880, label, 8, GRAY, align=PP_ALIGN.CENTER)
        stat_x += stat_w + 60000

    # 2张截图
    shots = [
        (ensure_cropped("admin_dashboard.png"), "管理员首页"),
        (ensure_cropped("admin_approvals.png"), "注册审批"),
        (ensure_cropped("admin_users.png"), "用户管理"),
    ]
    valid_shots = [(p, c) for p, c in shots if p]
    add_screenshot_row(slide, valid_shots, top=3350000, max_h=1350000,
                       total_width=8092440, left_margin=594360)
    add_textbox(slide, 594360, 4780000, 8092440, 182880,
                "来源：智启题库系统实拍（localhost:5173）", 8, GRAY_LIGHT, align=PP_ALIGN.CENTER)
    print(f"  管理员端: {len(valid_shots)} 张截图")

SLIDE_W = 9144000
SLIDE_H = 5143500

# ==================== 主逻辑 ====================
def main():
    prs = Presentation(PPTX_PATH)
    print(f"加载PPT: {len(prs.slides)} 张幻灯片")

    print("\n=== 重建幻灯片 11: 学生端 ===")
    rebuild_student_slide(prs.slides[10])

    print("\n=== 重建幻灯片 12: 教师端 ===")
    rebuild_teacher_slide(prs.slides[11])

    print("\n=== 重建幻灯片 13: 管理员端 ===")
    rebuild_admin_slide(prs.slides[12])

    prs.save(PPTX_PATH)
    print(f"\n✅ PPT 已保存: {PPTX_PATH}")

if __name__ == '__main__':
    main()
