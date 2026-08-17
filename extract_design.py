from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor
from lxml import etree

pptx_path = r'c:\Users\gaoxu\.trae-cn\attachments\6a7dee59db35f9bcbe161e89\b2f215b7-2b04-4adc-9ee0-86ae4abf6090_8b792b31-2b55-4c92-a079-e43ae7d69933_智启题库-大学生创新创业大赛答辩-优化版-修正.pptx'
prs = Presentation(pptx_path)

# Extract color and font info from slide 3 (a typical content slide)
slide = prs.slides[2]  # Slide 3

print("=== Slide 3 Design Details ===")
for shape in slide.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                font = run.font
                color_info = "default"
                try:
                    if font.color and font.color.rgb:
                        color_info = str(font.color.rgb)
                except:
                    color_info = "inherit"

                fill_info = "none"
                try:
                    if shape.fill and shape.fill.fore_color:
                        fill_info = str(shape.fill.fore_color.rgb)
                except:
                    fill_info = "inherit"

                print(f'  Text="{run.text[:50]}" font={font.name} size={font.size} bold={font.bold} color={color_info} fill={fill_info}')

# Check the slide layout XML for color scheme
print("\n=== Theme Colors ===")
theme = prs.slide_layouts[0]
xml = etree.tostring(theme.element, pretty_print=True).decode()
# Just print the first 2000 chars to see the structure
print(xml[:2000])

# Check slide background
print("\n=== Slide Background ===")
for slide in prs.slides:
    bg = slide.background
    print(f"Slide bg: {bg.fill.type if bg.fill else 'none'}")
    break
