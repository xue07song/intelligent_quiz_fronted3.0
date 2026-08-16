from pptx import Presentation

pptx_path = r'C:\Users\gaoxu\Desktop\智启题库-大学生创新创业大赛答辩-商业计划书.pptx'
prs = Presentation(pptx_path)

for idx in [10, 11, 12]:
    slide = prs.slides[idx]
    print(f'=== Slide {idx+1} ===')
    pic_count = 0
    text_count = 0
    for shape in slide.shapes:
        if shape.shape_type == 13:
            pic_count += 1
            print(f'  Picture: left={shape.left/914400:.2f}", top={shape.top/914400:.2f}", w={shape.width/914400:.2f}", h={shape.height/914400:.2f}"')
        elif shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                text_count += 1
                if len(t) < 60:
                    print(f'  Text: "{t}" at top={shape.top/914400:.2f}"')
    print(f'  Total: {pic_count} pictures, {text_count} text shapes')
    print()
