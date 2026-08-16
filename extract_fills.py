from pptx import Presentation
from pptx.util import Emu, Pt
from lxml import etree
import re

pptx_path = r'c:\Users\gaoxu\.trae-cn\attachments\6a7dee59db35f9bcbe161e89\b2f215b7-2b04-4adc-9ee0-86ae4abf6090_8b792b31-2b55-4c92-a079-e43ae7d69933_智启题库-大学生创新创业大赛答辩-优化版-修正.pptx'
prs = Presentation(pptx_path)

ns = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
}

# Extract fill colors from Slide 3 (typical content slide)
slide = prs.slides[2]
print("=== Slide 3 Shape Fills ===")
for i, shape in enumerate(slide.shapes):
    xml = etree.tostring(shape.element, pretty_print=False).decode()

    # Extract solidFill color
    solid_colors = re.findall(r'<a:srgbClr val="([0-9A-Fa-f]{6})"', xml)
    scheme_colors = re.findall(r'<a:schemeClr val="(\w+)"', xml)

    text = ""
    if shape.has_text_frame:
        text = shape.text_frame.text[:40]

    fill_type = "none"
    try:
        if shape.fill.type is not None:
            fill_type = str(shape.fill.type)
    except:
        pass

    print(f"Shape {i}: name={shape.name}, text='{text}', fill_type={fill_type}, solid_colors={solid_colors}, scheme_colors={scheme_colors}")

# Also check Slide 10 (student demo) for image placeholders
print("\n=== Slide 10 (Student Demo) Shapes ===")
slide10 = prs.slides[9]
for i, shape in enumerate(slide10.shapes):
    text = ""
    if shape.has_text_frame:
        text = shape.text_frame.text[:60]
    print(f"Shape {i}: name={shape.name}, type={shape.shape_type}, pos=({shape.left},{shape.top}), size=({shape.width},{shape.height}), text='{text}'")

# Check Slide 13 (competitive table) for table structure
print("\n=== Slide 13 (Competitive) Shapes ===")
slide13 = prs.slides[12]
for i, shape in enumerate(slide13.shapes):
    text = ""
    if shape.has_text_frame:
        text = shape.text_frame.text[:80]
    has_table = shape.has_table if hasattr(shape, 'has_table') else False
    print(f"Shape {i}: name={shape.name}, type={shape.shape_type}, text='{text}', has_table={has_table}")
