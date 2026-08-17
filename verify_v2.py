from pptx import Presentation

pptx_path = r'C:\Users\gaoxu\Desktop\智启题库-大学生创新创业大赛答辩-商业计划书.pptx'
prs = Presentation(pptx_path)

for idx in [10, 11, 12]:
    slide = prs.slides[idx]
    print(f'=== Slide {idx+1} ===')
    for shape in slide.shapes:
        if shape.shape_type == 13:
            print(f'  PIC: left={shape.left/914400:.2f}", top={shape.top/914400:.2f}", w={shape.width/914400:.2f}", h={shape.height/914400:.2f}"')
        elif shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t and len(t) < 50 and shape.top and shape.top < 1000000:
                print(f'  HDR: "{t}" top={shape.top/914400:.2f}"')
    print()
