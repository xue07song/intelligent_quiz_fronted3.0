from pptx import Presentation
from pptx.util import Emu

pptx_path = r'C:\Users\gaoxu\Desktop\智启题库-大学生创新创业大赛答辩-商业计划书.pptx'
prs = Presentation(pptx_path)

# Fix screenshot positions on demo slides (slides 11, 12, 13 = index 10, 11, 12)
# The screenshots were placed at top=4251960 with height=1280160, which extends beyond the slide
# We need to reposition them to fit within the available space

# Available space: cards end at ~4389120, page number at ~4823460
# Let's reduce card heights slightly and place screenshots below

for slide_idx in [10, 11, 12]:  # Demo slides (student, teacher, admin)
    slide = prs.slides[slide_idx]

    # Find all pictures in this slide
    pictures = []
    captions = []
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            pictures.append(shape)
        if shape.has_text_frame and '来源' in shape.text_frame.text:
            captions.append(shape)

    if not pictures:
        print(f"Slide {slide_idx+1}: No pictures found")
        continue

    # First, shrink the bottom row of cards to make room
    # Cards in bottom row start at top=2788920, height=1600200
    # Let's shrink them to height=1300000 (save 300200 EMU)
    for shape in slide.shapes:
        if shape.has_text_frame:
            continue
        if shape.shape_type == 1:  # AUTO_SHAPE (card backgrounds, accent bars)
            if shape.top and shape.top >= 2700000 and shape.top <= 2900000:
                # This is a bottom-row card element
                old_h = shape.height
                shape.height = Emu(1300000)
                print(f"  Slide {slide_idx+1}: Resized shape from h={old_h} to h=1300000")

    # Reposition screenshots to fit in the new space
    # New available area: top=4200000, height=580000 (0.63 inches)
    shot_y = 4150000
    shot_h = 580000
    shot_w = 2600000  # 2.84 inches
    shot_x_start = 594360
    shot_gap = 274320

    for i, pic in enumerate(pictures):
        new_x = shot_x_start + i * (shot_w + shot_gap)
        pic.left = Emu(new_x)
        pic.top = Emu(shot_y)
        pic.width = Emu(shot_w)
        pic.height = Emu(shot_h)
        print(f"  Slide {slide_idx+1}: Repositioned picture {i+1} to ({new_x}, {shot_y}), size=({shot_w}, {shot_h})")

    # Reposition captions
    for i, cap in enumerate(captions):
        new_x = shot_x_start + i * (shot_w + shot_gap)
        cap.left = Emu(new_x)
        cap.top = Emu(shot_y + shot_h - 50000)
        cap.width = Emu(shot_w)
        print(f"  Slide {slide_idx+1}: Repositioned caption {i+1}")

prs.save(pptx_path)
print(f"\n✅ Screenshot positions fixed and saved")
