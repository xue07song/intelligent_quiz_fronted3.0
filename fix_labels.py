from pptx import Presentation

pptx_path = r'C:\Users\gaoxu\Desktop\智启题库-大学生创新创业大赛答辩-商业计划书.pptx'
prs = Presentation(pptx_path)

# Fix section labels that weren't updated correctly
# We need to find the actual text in each slide and replace it
fixes = [
    (15, None, "08  PROMOTION"),   # Slide 16 (index 15): 应用推广路径
    (17, None, "10  TEAM"),         # Slide 18 (index 17): 团队介绍
    (19, None, "12  ROADMAP"),      # Slide 20 (index 19): 未来发展规划
]

for slide_idx, _, new_text in fixes:
    slide = prs.slides[slide_idx]
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            # Check if this looks like a section label (contains uppercase English)
            if text and any(c.isupper() for c in text) and len(text) < 40:
                # Check if it's in the section label position (top area)
                if shape.top and shape.top < 1000000 and shape.top > 600000:
                    print(f"  Slide {slide_idx+1}: Found label '{text}' at top={shape.top}")
                    for p in shape.text_frame.paragraphs:
                        for r in p.runs:
                            r.text = new_text
                    print(f"    → Updated to '{new_text}'")

# Also verify all section labels
print("\n=== All Section Labels ===")
for i, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text and any(c.isupper() for c in text) and len(text) < 40:
                if shape.top and shape.top < 1000000 and shape.top > 600000:
                    print(f"  Slide {i+1}: '{text}'")

prs.save(pptx_path)
print(f"\n✅ Labels fixed and saved")
