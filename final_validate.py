from pptx import Presentation

pptx_path = r'C:\Users\gaoxu\Desktop\智启题库-大学生创新创业大赛答辩-商业计划书.pptx'
prs = Presentation(pptx_path)
print(f'Total slides: {len(prs.slides)}')
print()

modules = [
    "01 项目背景", "02 市场调研", "03 项目定位", "04 技术原理",
    "05 产品功能", "06 竞品对比", "07 商业模式", "08 推广规划",
    "09 知识产权", "10 团队介绍", "11 社会价值", "12 未来规划", "13 总结展望"
]

for i, slide in enumerate(prs.slides):
    texts = []
    pic_count = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t and len(t) < 80:
                texts.append(t)
        if shape.shape_type == 13:
            pic_count += 1
    title = texts[0] if texts else '(no text)'
    label = texts[1] if len(texts) > 1 else ''
    print(f'  {i+1:2d}. [{label:20s}] {title[:40]:40s} pics={pic_count}')

print()
# Check 13 modules coverage
labels_found = []
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t and len(t) < 30:
                labels_found.append(t)

print("Module coverage check:")
for m in modules:
    found = any(m.split()[1] in l for l in labels_found)
    status = "✓" if found else "✗"
    print(f"  {status} {m}")

# Check for watermarks
watermark_count = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if 'AI' in t and '生成' in t:
                watermark_count += 1
print(f"\nWatermarks found: {watermark_count}")
