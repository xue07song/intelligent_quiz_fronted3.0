from pptx import Presentation
from pptx.util import Emu
import os

pptx_path = r'C:\Users\gaoxu\Desktop\智启题库-大学生创新创业大赛答辩-商业计划书.pptx'
prs = Presentation(pptx_path)

for idx in [10, 11, 12]:  # slides 11, 12, 13
    slide = prs.slides[idx]
    print(f'=== Slide {idx+1} ===')
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            print(f'  Picture: left={shape.left}, top={shape.top}, w={shape.width}, h={shape.height}')
            print(f'    Inches: left={shape.left/914400:.2f}, top={shape.top/914400:.2f}, w={shape.width/914400:.2f}, h={shape.height/914400:.2f}')
            try:
                img = shape.image
                print(f'    Image blob size: {len(img.blob)} bytes, ext: {img.ext}')
            except:
                pass
        elif shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                print(f'  Text: "{t[:80]}" at top={shape.top}')
    print()
