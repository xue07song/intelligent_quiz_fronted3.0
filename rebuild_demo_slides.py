"""
重建功能展示幻灯片 (11/12/13)：
- 压缩文字内容，放大截图
- 替换学生端截图为更清晰的版本
- 所有截图保持原始宽高比
"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image as PILImage

# ==================== 常量 ====================
PPTX_PATH = r'C:\Users\gaoxu\Desktop\智启题库-大学生创新创业大赛答辩-商业计划书.pptx'
SCREENSHOT_DIR = r'c:\Users\gaoxu\AppData\Local\Temp\trae\screenshots'

NAVY = RGBColor(0x1A, 0x2B, 0x5E)
TEAL = RGBColor(0x00, 0xB8, 0xA3)
BLUE = RGBColor(0x2D, 0x6C, 0xDF)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT = RGBColor(0x9C, 0xA3, 0xAF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT_NAME = "Microsoft YaHei"

SLIDE_W = 9144000
SLIDE_H = 5143500
SIDEBAR_W = 73152

# ==================== 辅助函数 ====================
def set_font(run, size_pt, color, bold=False, name=FONT_NAME):
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = name

def add_textbox(slide, left, top, width, height, text, size_pt, color, bold=False, align=None):
    box = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_font(run, size_pt, color, bold)
    if align:
        p.alignment = align
    return box

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_compact_card(slide, left, top, width, height, accent_color, title, items, title_size=13):
    """紧凑型内容卡片：标题 + 2条要点"""
    add_rect(slide, left, top, width, height, WHITE)
    add_rect(slide, left, top, 45720, height, accent_color)
    add_textbox(slide, left + 137160, top + 91440, width - 228600, 228600, title, title_size, NAVY, True)
    y = top + 365760
    for item in items:
        add_textbox(slide, left + 137160, y, width - 228600, 201168, f"•  {item}", 9, GRAY)
        y += 256032
    return y

def add_tall_card(slide, left, top, width, height, accent_color, title, items, title_size=13):
    """较高内容卡片：标题 + 3-4条要点"""
    add_rect(slide, left, top, width, height, WHITE)
    add_rect(slide, left, top, 45720, height, accent_color)
    add_textbox(slide, left + 137160, top + 91440, width - 228600, 228600, title, title_size, NAVY, True)
    y = top + 365760
    for item in items:
        add_textbox(slide, left + 137160, y, width - 228600, 201168, f"•  {item}", 9, GRAY)
        y += 256032
    return y

def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)

def clear_slide_content(slide, keep_count=4):
    """保留前 keep_count 个 shape（sidebar, title, label, page），删除其余"""
    shapes_to_remove = []
    for i, shape in enumerate(slide.shapes):
        if i >= keep_count:
            shapes_to_remove.append(shape)
    for s in shapes_to_remove:
        remove_shape(s)

def get_image_aspect_ratio(path):
    """获取图片宽高比 (width / height)"""
    with PILImage.open(path) as img:
        w, h = img.size
        return w / h

def fit_image_in_area(path, max_w, max_h):
    """在给定区域内计算图片尺寸，保持宽高比"""
    ratio = get_image_aspect_ratio(path)
    # 先按宽度适配
    w = max_w
    h = int(w / ratio)
    # 如果高度超出，按高度适配
    if h > max_h:
        h = max_h
        w = int(h * ratio)
    return w, h

def add_screenshot_with_label(slide, img_path, left, top, max_w, max_h, caption="来源：智启题库系统实拍"):
    """添加截图（保持宽高比）+ 来源标注"""
    w, h = fit_image_in_area(img_path, max_w, max_h)
    # 居中对齐
    x = left + (max_w - w) // 2
    y = top + (max_h - h) // 2
    pic = slide.shapes.add_picture(img_path, Emu(x), Emu(y), Emu(w), Emu(h))
    # 来源标注
    add_textbox(slide, left, top + max_h + 20000, max_w, 182880, caption, 7, GRAY_LIGHT, align=PP_ALIGN.CENTER)
    return pic

def find_screenshot(name):
    p = os.path.join(SCREENSHOT_DIR, name)
    if os.path.exists(p):
        return p
    return None

# ==================== 重建函数 ====================
def rebuild_student_slide(slide):
    """重建学生端功能展示幻灯片"""
    clear_slide_content(slide, keep_count=4)

    # 副标题
    add_textbox(slide, 594360, 1005840, 8092440, 274320,
                "6种题型在线答题 · 自适应难度练习 · 错题归集分析 · AI智能助手", 11, GRAY)

    # 4张紧凑卡片 (2x2)
    cw = 3977640  # 卡片宽度
    ch = 820000   # 卡片高度
    cx1 = 457200
    cx2 = 457200 + cw + 228600  # 4659840
    cy1 = 1371600
    cy2 = 1371600 + ch + 160000  # 2351600

    add_compact_card(slide, cx1, cy1, cw, ch, TEAL, "在线答题", [
        "6种题型：判断/单选/多选/填空/简答/程序",
        "答题卡导航 · 实时计时 · 断点续答",
    ])
    add_compact_card(slide, cx2, cy1, cw, ch, BLUE, "自适应练习", [
        "从1级开始，逐题动态调整难度",
        "库存预检 + 方案推荐 · AI即时答疑",
    ])
    add_compact_card(slide, cx1, cy2, cw, ch, AMBER, "错题本与分析", [
        "错题自动归集 · 一键重练",
        "章节/题型/难度三维统计 · SVG趋势图",
    ])
    add_compact_card(slide, cx2, cy2, cw, ch, NAVY, "AI智能助手", [
        "浮动球设计 · 可拖拽对话面板",
        "自然语言组卷 · 错题浓缩 · 同类题推荐",
    ])

    # 截图区域
    shot_top = 3250000
    shot_max_h = 1500000
    shot_max_w = 3900000
    shot_left1 = 594360
    shot_left2 = 594360 + shot_max_w + 294000  # 4788360

    for fname, left in [("student_exam.png", shot_left1), ("student_practice.png", shot_left2)]:
        path = find_screenshot(fname)
        if path:
            add_screenshot_with_label(slide, path, left, shot_top, shot_max_w, shot_max_h)
            print(f"  插入截图: {fname}")
        else:
            print(f"  截图未找到: {fname}")

def rebuild_teacher_slide(slide):
    """重建教师端功能展示幻灯片"""
    clear_slide_content(slide, keep_count=4)

    # 副标题
    add_textbox(slide, 594360, 1005840, 8092440, 274320,
                "多约束智能组卷 · AI自动出题 · 试卷分析与学情总览", 11, GRAY)

    # 4张紧凑卡片 (2x2)
    cw = 3977640
    ch = 820000
    cx1 = 457200
    cx2 = 4659840
    cy1 = 1371600
    cy2 = 2351600

    add_compact_card(slide, cx1, cy1, cw, ch, TEAL, "多约束智能组卷", [
        "四步配置 · 9套预设 + 手动自定义 + AI辅助",
        "题型×难度×章节×知识点 四维校验 · 替代方案",
    ])
    add_compact_card(slide, cx2, cy1, cw, ch, BLUE, "AI自动出题", [
        "按章节/知识点/题型/难度生成草稿",
        "人工审核 → 入库 · Excel批量导入 · 6种题型",
    ])
    add_compact_card(slide, cx1, cy2, cw, ch, AMBER, "试卷分析与学情", [
        "正确率分析 · 成绩分布 · 班级对比",
        "主观题复核 · 学生个性化分析",
    ])
    add_compact_card(slide, cx2, cy2, cw, ch, NAVY, "自适应学情总览", [
        "练习人数 · 有效次数 · 整体正确率",
        "关注学生标记 · 建议跳转组卷",
    ])

    # 截图区域
    shot_top = 3250000
    shot_max_h = 1500000
    shot_max_w = 3900000
    shot_left1 = 594360
    shot_left2 = 4788360

    for fname, left in [("teacher_exam_gen.png", shot_left1), ("teacher_analytics.png", shot_left2)]:
        path = find_screenshot(fname)
        if path:
            add_screenshot_with_label(slide, path, left, shot_top, shot_max_w, shot_max_h)
            print(f"  插入截图: {fname}")
        else:
            print(f"  截图未找到: {fname}")

def rebuild_admin_slide(slide):
    """重建管理员端功能展示幻灯片"""
    clear_slide_content(slide, keep_count=4)

    # 副标题
    add_textbox(slide, 594360, 1005840, 8092440, 274320,
                "三角色RBAC权限管理 · 注册审批闭环 · 全局数据治理", 11, GRAY)

    # 3张内容卡片
    cw = 2545600
    ch = 1450000
    cx1 = 457200
    cx2 = 457200 + cw + 228600   # 3231400
    cx3 = 3231400 + cw + 228600  # 6005600
    cy = 1371600

    add_tall_card(slide, cx1, cy, cw, ch, TEAL, "用户与权限管理", [
        "三角色RBAC：admin/teacher/student",
        "教师科目级数据隔离",
        "用户启停 · 重置密码 · 批量操作",
        "JWT 7天有效 · bcrypt加盐哈希",
    ])
    add_tall_card(slide, cx2, cy, cw, ch, BLUE, "注册审批闭环", [
        "非自助注册 · 审批制",
        "教师申请含学院/专业/科目/工号",
        "管理员审核 → 通过/拒绝",
        "通过后自动创建账号并关联科目",
    ])
    add_tall_card(slide, cx3, cy, cw, ch, AMBER, "全局数据治理", [
        "全局题库与学习数据统计",
        "增量字段兼容 · 减少迁移风险",
        "用户反馈管理：提交/回复/跟踪",
        "系统运行监控与日志",
    ])

    # 数据统计指标行
    stats = [("70+", "API接口"), ("15+", "数据表"), ("30+", "前端组件"), ("6", "题型支持"), ("2", "大模型协作")]
    stat_y = 2950000
    stat_w = 1500000
    stat_gap = (8092440 - 5 * stat_w) // 4
    stat_x = 594360
    for num, label in stats:
        add_textbox(slide, stat_x, stat_y, stat_w, 274320, num, 18, TEAL, True, align=PP_ALIGN.CENTER)
        add_textbox(slide, stat_x, stat_y + 297176, stat_w, 201168, label, 9, GRAY, align=PP_ALIGN.CENTER)
        stat_x += stat_w + stat_gap

    # 截图区域
    shot_top = 3650000
    shot_max_h = 1050000
    shot_max_w = 3900000
    shot_left1 = 594360
    shot_left2 = 4788360

    for fname, left in [("admin_dashboard.png", shot_left1), ("admin_approvals.png", shot_left2)]:
        path = find_screenshot(fname)
        if path:
            add_screenshot_with_label(slide, path, left, shot_top, shot_max_w, shot_max_h)
            print(f"  插入截图: {fname}")
        else:
            print(f"  截图未找到: {fname}")

# ==================== 主逻辑 ====================
def main():
    prs = Presentation(PPTX_PATH)
    print(f"加载PPT: {len(prs.slides)} 张幻灯片")

    # Slide 11 (index 10): 学生端
    print("\n=== 重建幻灯片 11: 学生端 ===")
    rebuild_student_slide(prs.slides[10])

    # Slide 12 (index 11): 教师端
    print("\n=== 重建幻灯片 12: 教师端 ===")
    rebuild_teacher_slide(prs.slides[11])

    # Slide 13 (index 12): 管理员端
    print("\n=== 重建幻灯片 13: 管理员端 ===")
    rebuild_admin_slide(prs.slides[12])

    # 保存
    prs.save(PPTX_PATH)
    print(f"\n✅ PPT 已保存: {PPTX_PATH}")

if __name__ == '__main__':
    main()
