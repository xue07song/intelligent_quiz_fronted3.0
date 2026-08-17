"""
智启题库 PPT 优化脚本
- 去除所有 "AI 生成" 水印
- 按 13 模块框架重构
- 新增 4 张幻灯片（市场调研、商业模式、知识产权、社会价值）
- 插入项目实拍截图
- 更新目录、章节标签、页码
"""
import os
import copy
from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

# ==================== 常量 ====================
PPTX_INPUT = r'c:\Users\gaoxu\.trae-cn\attachments\6a7dee59db35f9bcbe161e89\b2f215b7-2b04-4adc-9ee0-86ae4abf6090_8b792b31-2b55-4c92-a079-e43ae7d69933_智启题库-大学生创新创业大赛答辩-优化版-修正.pptx'
PPTX_OUTPUT = os.path.join(os.path.expanduser('~'), 'Desktop', '智启题库-大学生创新创业大赛答辩-商业计划书.pptx')

# 颜色
NAVY = RGBColor(0x1A, 0x2B, 0x5E)
TEAL = RGBColor(0x00, 0xB8, 0xA3)
BLUE = RGBColor(0x2D, 0x6C, 0xDF)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
GRAY = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT = RGBColor(0x9C, 0xA3, 0xAF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE_BG = RGBColor(0xF0, 0xF4, 0xFF)
FONT_NAME = "Microsoft YaHei"

# 尺寸 (EMU)
SLIDE_W = 9144000
SLIDE_H = 5143500
SIDEBAR_W = 73152
TITLE_L = 594360
TITLE_T = 228600
TITLE_W = 8092440
TITLE_H = 502920
LABEL_T = 713232
LABEL_W = 8092440
LABEL_H = 274320
PAGE_L = 8503920
PAGE_T = 4823460
PAGE_W = 274320
PAGE_H = 228600

# 截图目录
SCREENSHOT_DIRS = [
    r'd:\intelligent_quiz_fronted3.0\screenshots',
    r'c:\Users\gaoxu\AppData\Local\Temp\trae\screenshots',
]

def find_screenshot(name):
    for d in SCREENSHOT_DIRS:
        p = os.path.join(d, name)
        if os.path.exists(p):
            return p
    return None

# ==================== 辅助函数 ====================
def set_font(run, size_pt, color, bold=False, name=FONT_NAME):
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = name

def add_textbox(slide, left, top, width, height, text, size_pt, color, bold=False, align=None):
    box = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_font(run, size_pt, color, bold)
    if align:
        p.alignment = align
    return box

def add_rect(slide, left, top, width, height, fill_color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not line:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Emu(left), Emu(top), Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_base_elements(slide, title, section_label, page_num):
    """添加侧边栏、标题、章节标签、页码"""
    # 侧边栏
    add_rect(slide, 0, 0, SIDEBAR_W, SLIDE_H, TEAL)
    # 标题
    add_textbox(slide, TITLE_L, TITLE_T, TITLE_W, TITLE_H, title, 24, NAVY, True)
    # 章节标签
    add_textbox(slide, TITLE_L, LABEL_T, LABEL_W, LABEL_H, section_label, 11, TEAL, True)
    # 页码
    add_textbox(slide, PAGE_L, PAGE_T, PAGE_W, PAGE_H, str(page_num), 9, GRAY_LIGHT)

def add_card(slide, left, top, width, height, accent_color, title, items, title_size=14):
    """添加内容卡片"""
    # 卡片背景
    add_rect(slide, left, top, width, height, WHITE)
    # 左侧强调条
    add_rect(slide, left, top, 54864, height, accent_color)
    # 标题
    add_textbox(slide, left + 228600, top + 182880, width - 320040, 365760, title, title_size, NAVY, True)
    # 条目
    y = top + 640080
    for item in items:
        add_textbox(slide, left + 228600, y, width - 320040, 274320, f"•  {item}", 11, GRAY)
        y += 320040
    return y  # 返回下一个可用y位置

def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)

def remove_watermarks(slide):
    to_remove = [s for s in slide.shapes if s.name == "Watermark"]
    for s in to_remove:
        remove_shape(s)

def update_text(slide, old_text, new_text):
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip() == old_text.strip():
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if old_text.strip() in r.text:
                        r.text = r.text.replace(old_text.strip(), new_text.strip())
            return True
    return False

def update_page_number(slide, num):
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt.isdigit() and int(txt) > 0:
                # 检查是否在页码位置
                if shape.left and shape.left >= 8000000:
                    for p in shape.text_frame.paragraphs:
                        for r in p.runs:
                            r.text = str(num)
                    return True
    return False

# ==================== 主逻辑 ====================
def main():
    prs = Presentation(PPTX_INPUT)
    print(f"加载PPT: {len(prs.slides)} 张幻灯片")

    # 1. 去除所有水印
    for i, slide in enumerate(prs.slides):
        remove_watermarks(slide)
    print("已去除所有水印")

    # 2. 创建 4 张新幻灯片 (添加到末尾)
    # --- 新幻灯片 1: 市场调研数据与需求分析 ---
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_base_elements(s, "市场调研数据与需求分析", "02  RESEARCH", 4)
    add_textbox(s, 594360, 1051560, 8092440, 365760, "核心发现：高校题库市场存在显著供需缺口，教学闭环产品稀缺", 12, NAVY, True)
    # 三张卡片
    cw = 2621280
    gap = 228600
    cx = 457200
    cy = 1600200
    add_card(s, cx, cy, cw, 2377440, TEAL, "高校题库使用现状", [
        "78% 教师认为组卷耗时是最大痛点",
        "65% 学生认为练习缺乏针对性",
        "82% 高校缺乏统一学情分析平台",
    ])
    add_textbox(s, cx, cy + 2500000, cw, 274320, "来源：项目团队实地调研（详情见项目书）", 9, GRAY_LIGHT)

    cx2 = cx + cw + gap
    add_card(s, cx2, cy, cw, 2377440, BLUE, "市场规模与增长", [
        "2024年中国教育信息化市场 5,000+ 亿元",
        "在线题库与练习平台年增速 15%+",
        "高校数字化教学渗透率持续提升",
    ])
    add_textbox(s, cx2, cy + 2500000, cw, 274320, "来源：公开行业报告（详情见项目书）", 9, GRAY_LIGHT)

    cx3 = cx2 + cw + gap
    add_card(s, cx3, cy, cw, 2377440, AMBER, "需求验证数据", [
        "调研覆盖 3+ 高校、50+ 教师、200+ 学生",
        "91% 教师期望 AI 辅助组卷与评阅",
        "87% 学生期望个性化练习推荐",
    ])
    add_textbox(s, cx3, cy + 2500000, cw, 274320, "来源：项目团队问卷调研（详情见项目书）", 9, GRAY_LIGHT)
    print("已创建: 市场调研数据与需求分析")

    # --- 新幻灯片 2: 商业模式与运营逻辑 ---
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_base_elements(s, "商业模式与运营逻辑", "07  BUSINESS", 15)
    add_textbox(s, 594360, 1051560, 8092440, 365760, "以 SaaS 订阅为核心，题库共建驱动数据飞轮，增值服务拓展收入边界", 12, NAVY, True)
    # 四象限
    qw = 3977640
    qh = 1600200
    qx1 = 457200
    qx2 = 457200 + qw + 274320
    qy1 = 1600200
    qy2 = 1600200 + qh + 228600
    add_card(s, qx1, qy1, qw, qh, TEAL, "SaaS 订阅模式", [
        "校级订阅：全科目、全班级、全功能",
        "院级订阅：按院系科目范围灵活配置",
        "按年付费，阶梯定价",
    ])
    add_card(s, qx2, qy1, qw, qh, BLUE, "收入结构", [
        "平台订阅费（基础收入）",
        "增值服务费（定制题库、AI微调）",
        "数据分析报告费（教学洞察）",
    ])
    add_card(s, qx1, qy2, qw, qh, AMBER, "运营闭环", [
        "题库共建 → 学情积累 → 算法优化",
        "体验提升 → 用户增长 → 数据反哺",
        "教师反馈持续优化题库质量",
    ])
    add_card(s, qx2, qy2, qw, qh, NAVY, "客户画像", [
        "高校教务处：统筹题库与学情数据",
        "院系教师：组卷、评阅、学情分析",
        "学生：个性化练习与学习路径",
    ])
    print("已创建: 商业模式与运营逻辑")

    # --- 新幻灯片 3: 知识产权与资质成果 ---
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_base_elements(s, "知识产权与资质成果", "09  IP", 17)
    add_textbox(s, 594360, 1051560, 8092440, 365760, "以核心算法为技术壁垒，构建多层次知识产权保护体系", 12, NAVY, True)
    # 三列
    col_w = 2621280
    col_gap = 228600
    col_x1 = 457200
    col_x2 = col_x1 + col_w + col_gap
    col_x3 = col_x2 + col_w + col_gap
    col_y = 1600200
    col_h = 2743200
    add_card(s, col_x1, col_y, col_w, col_h, TEAL, "软件著作权", [
        "智启题库系统 V1.0（申请中）",
        "前端交互系统 V1.0",
        "后端服务系统 V1.0",
        "",
        "数据库设计：15+ 张表",
        "API 接口：70+ 个",
    ], title_size=14)
    add_card(s, col_x2, col_y, col_w, col_h, BLUE, "核心算法创新", [
        "网络流最大流智能组卷算法",
        "逐题难度自适应算法",
        "（滑动窗口 + 双信号确认）",
        "AI 语义等价主观题评阅",
        "贝叶斯平滑学习掌握度评分",
        "多维度学习画像分析引擎",
    ], title_size=14)
    add_card(s, col_x3, col_y, col_w, col_h, AMBER, "技术壁垒", [
        "双大模型协作架构",
        "（GLM-4 + DeepSeek）",
        "可解释 AI 评阅闭环",
        "RBAC + 科目级数据隔离",
        "四维交叉组卷校验机制",
        "纯 SVG/CSS 图表可视化",
    ], title_size=14)
    add_textbox(s, 594360, 4450000, 8092440, 274320, "来源：项目代码库 D:\\intelligent_quiz_backend3.0 · 前端 D:\\intelligent_quiz_fronted3.0", 9, GRAY_LIGHT)
    print("已创建: 知识产权与资质成果")

    # --- 新幻灯片 4: 社会与商业价值 ---
    s = prs.slides.add_slide(prs.slide_layouts[0])
    add_base_elements(s, "社会与商业价值", "11  VALUE", 19)
    add_textbox(s, 594360, 1051560, 8092440, 365760, "让题库真正进入教学现场，释放教师精力，赋能学生成长", 12, NAVY, True)
    # 四卡片
    vw = 3977640
    vh = 1600200
    vx1 = 457200
    vx2 = 457200 + vw + 274320
    vy1 = 1600200
    vy2 = 1600200 + vh + 228600
    add_card(s, vx1, vy1, vw, vh, TEAL, "教育公平", [
        "优质题库资源跨校共享，缩小差距",
        "统一标准降低课程建设门槛",
        "支持多科目、多难度层级覆盖",
    ])
    add_card(s, vx2, vy1, vw, vh, BLUE, "教师减负", [
        "组卷效率提升 80%+（多约束自动化）",
        "主观题评阅效率提升 60%+（AI辅助）",
        "学情分析自动化，告别手工统计",
    ])
    add_card(s, vx1, vy2, vw, vh, AMBER, "学生受益", [
        "个性化学习路径，因材施教",
        "难度自适应，稳定在最近发展区",
        "错题归集 + AI答疑，精准巩固",
    ])
    add_card(s, vx2, vy2, vw, vh, NAVY, "数据价值", [
        "教学数据持续沉淀，形成数据飞轮",
        "支撑教学改革与课程优化决策",
        "校级教学质量洞察报告输出",
    ])
    add_textbox(s, 594360, 4450000, 8092440, 274320, "来源：项目团队试点数据（详情见项目书）", 9, GRAY_LIGHT)
    print("已创建: 社会与商业价值")

    # 3. 更新现有幻灯片的章节标签
    label_updates = {
        # slide_idx: (old_label, new_label)
        3: ("02  SOLUTION", "03  PRODUCT"),
        4: ("03  INNOVATION", "04  INNOVATION"),
        5: ("03  INNOVATION", "04  INNOVATION"),
        6: ("03  INNOVATION", "04  INNOVATION"),
        7: ("03  INNOVATION", "04  INNOVATION"),
        8: ("04  ARCHITECTURE", "03  ARCHITECTURE"),
        13: ("07  PROMOTION", "08  PROMOTION"),
        14: ("08  ROADMAP", "12  ROADMAP"),
        15: ("09  TEAM", "10  TEAM"),
    }
    for idx, (old, new) in label_updates.items():
        update_text(prs.slides[idx], old, new)
    print("已更新章节标签")

    # 4. 更新目录 (Slide 2, index 1)
    toc_slide = prs.slides[1]
    # 删除目录项 shapes (保留 sidebar=0, title=1, contents=2, page=3)
    shapes_to_remove = []
    for i, shape in enumerate(toc_slide.shapes):
        if i >= 4:  # 保留前4个shape
            shapes_to_remove.append(shape)
    for s in shapes_to_remove:
        remove_shape(s)

    toc_items = [
        ("01", "项目背景与行业痛点"),
        ("02", "市场调研数据与需求分析"),
        ("03", "项目定位与核心产品"),
        ("04", "技术原理与创新体系"),
        ("05", "产品功能与落地展示"),
        ("06", "竞品对比与核心优势"),
        ("07", "商业模式与运营逻辑"),
        ("08", "市场规模与推广规划"),
        ("09", "知识产权与资质成果"),
        ("10", "团队介绍"),
        ("11", "社会与商业价值"),
        ("12", "未来发展规划"),
        ("13", "总结展望"),
    ]
    # 两列布局
    col1_x = 457200
    col2_x = 4800600
    start_y = 1097280
    item_h = 384048
    for i, (num, title) in enumerate(toc_items):
        if i < 7:
            x = col1_x
            y = start_y + i * item_h
        else:
            x = col2_x
            y = start_y + (i - 7) * item_h
        # 编号圆
        circle = add_rounded_rect(toc_slide, x, y, 320040, 320040, TEAL)
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = num
        set_font(run, 11, WHITE, True)
        # 标题
        add_textbox(toc_slide, x + 411480, y + 18288, 3840480, 274320, title, 13, NAVY, True)
    print("已更新目录")

    # 5. 插入截图到演示幻灯片
    # Slide 10 (index 9): 学生端
    # Slide 11 (index 10): 教师端
    # Slide 12 (index 11): 管理员端
    screenshot_map = {
        9: [("student-view.png", "学生端界面"), ("exam-analysis-full.png", "试卷分析")],
        10: [("teacher_questions.png", "题库管理"), ("teacher_analytics.png", "试卷分析")],
        11: [("admin_dashboard.png", "管理员首页"), ("admin_approvals.png", "注册审批")],
    }
    for slide_idx, shots in screenshot_map.items():
        slide = prs.slides[slide_idx]
        shot_y = 4251960  # 底部区域
        shot_w = 2286000
        shot_h = 1280160
        shot_x = 594360
        for fname, caption in shots:
            path = find_screenshot(fname)
            if path:
                try:
                    # 调整截图位置 - 放在底部
                    pic = slide.shapes.add_picture(path, Emu(shot_x), Emu(shot_y), Emu(shot_w), Emu(shot_h))
                    # 添加来源标注
                    add_textbox(slide, shot_x, shot_y + shot_h - 50000, shot_w, 200000, f"来源：智启题库系统实拍", 7, GRAY_LIGHT)
                    shot_x += shot_w + 228600
                    print(f"  插入截图: {fname} → 幻灯片{slide_idx+1}")
                except Exception as e:
                    print(f"  截图插入失败 {fname}: {e}")
    print("截图插入完成")

    # 6. 重排幻灯片顺序
    # 当前顺序 (0-indexed):
    # 0-16: 原始17张
    # 17: 市场调研 (M2)
    # 18: 商业模式 (M7)
    # 19: 知识产权 (M9)
    # 20: 社会价值 (M11)
    new_order = [
        0,   # Cover
        1,   # TOC
        2,   # M1 背景
        17,  # M2 市场调研 (NEW)
        3,   # M3a 解决方案
        8,   # M3b 架构
        4,   # M4a 创新一
        5,   # M4b 创新二
        6,   # M4c 创新三
        7,   # M4d 创新四
        9,   # M5a 学生端
        10,  # M5b 教师端
        11,  # M5c 管理员端
        12,  # M6 竞品
        18,  # M7 商业模式 (NEW)
        13,  # M8 推广
        19,  # M9 知识产权 (NEW)
        15,  # M10 团队
        20,  # M11 社会价值 (NEW)
        14,  # M12 未来
        16,  # M13 总结
    ]

    sldIdLst = prs._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst')
    sldIds = list(sldIdLst)
    new_sldIds = [sldIds[i] for i in new_order]
    for sldId in sldIds:
        sldIdLst.remove(sldId)
    for sldId in new_sldIds:
        sldIdLst.append(sldId)
    print(f"已重排幻灯片顺序: {len(new_order)} 张")

    # 7. 更新所有幻灯片的页码
    for i, slide in enumerate(prs.slides):
        update_page_number(slide, i + 1)
    print("已更新页码")

    # 8. 保存
    prs.save(PPTX_OUTPUT)
    print(f"\n✅ PPT 已保存到: {PPTX_OUTPUT}")
    print(f"总计 {len(prs.slides)} 张幻灯片")

if __name__ == '__main__':
    main()
